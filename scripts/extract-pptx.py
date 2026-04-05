#!/usr/bin/env python3
"""
extract-pptx.py — Extract slide content from a .pptx file to JSON.

Usage:
    python3 extract-pptx.py <file.pptx> [--output slides.json]

Output JSON structure:
{
  "title": "Deck title (from first slide or filename)",
  "slide_count": 7,
  "slides": [
    {
      "index": 1,
      "layout": "title",          // guessed layout type
      "heading": "Main Title",
      "subheading": "Subtitle",
      "body": ["bullet 1", "bullet 2"],
      "notes": "Speaker notes text",
      "images": ["image1.png"],   // filenames of embedded images
      "tables": [                 // list of tables, each a list of rows
        [["Cell A1", "Cell B1"], ["Cell A2", "Cell B2"]]
      ]
    },
    ...
  ]
}

Requires: python-pptx
    pip install python-pptx
"""

import sys
import json
import argparse
from pathlib import Path


def _text(shape) -> str:
    """Return all text from a shape, stripped."""
    try:
        return shape.text_frame.text.strip()
    except AttributeError:
        return ""


def _paragraphs(shape) -> list[str]:
    """Return non-empty paragraph strings from a text frame."""
    try:
        return [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
    except AttributeError:
        return []


def _guess_layout(slide, idx: int) -> str:
    """Heuristic: guess slide layout type from content."""
    try:
        layout_name = slide.slide_layout.name.lower()
    except Exception:
        layout_name = ""

    if idx == 0 or "title" in layout_name:
        return "title"
    if "blank" in layout_name:
        return "blank"
    if "two" in layout_name or "comparison" in layout_name:
        return "two-col"
    if "picture" in layout_name or "image" in layout_name:
        return "image-full"
    if "quote" in layout_name:
        return "quote"
    return "content"


def _extract_images(slide) -> list[str]:
    """Return filenames of picture shapes embedded in the slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                ext  = shape.image.ext
                name = f"slide_img_{shape.shape_id}.{ext}"
                images.append(name)
            except Exception:
                pass
    return images


def _extract_tables(slide) -> list[list[list[str]]]:
    """Return a list of tables; each table is a list of rows of cell strings."""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            tbl = []
            for row in shape.table.rows:
                tbl.append([cell.text.strip() for cell in row.cells])
            tables.append(tbl)
    return tables


def _extract_notes(slide) -> str:
    """Return speaker notes text."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        return tf.text.strip()
    except Exception:
        return ""


def extract(path: Path) -> dict:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        print("Error: python-pptx is not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(path))
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        heading    = ""
        subheading = ""
        body_lines: list[str] = []

        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        # Placeholder indices: 0=title, 1=body/content, 2=subtitle (on title slides)
        if 0 in placeholders:
            heading = _text(placeholders[0])
        if 2 in placeholders:
            subheading = _text(placeholders[2])
        if 1 in placeholders:
            lines = _paragraphs(placeholders[1])
            if lines:
                body_lines = lines

        # Collect text from non-placeholder text boxes
        extra_texts = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            t = _text(shape)
            if t:
                extra_texts.append(t)

        slides_data.append({
            "index":      idx + 1,
            "layout":     _guess_layout(slide, idx),
            "heading":    heading,
            "subheading": subheading,
            "body":       body_lines,
            "extra":      extra_texts,
            "notes":      _extract_notes(slide),
            "images":     _extract_images(slide),
            "tables":     _extract_tables(slide),
        })

    # Deck title: first slide heading or filename stem
    deck_title = slides_data[0]["heading"] if slides_data else path.stem

    return {
        "title":       deck_title,
        "slide_count": len(slides_data),
        "slides":      slides_data,
    }


def main():
    parser = argparse.ArgumentParser(description="Extract PPTX slide content to JSON")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--output", "-o", default="-",
                        help="Output file path (default: stdout)")
    args = parser.parse_args()

    path = Path(args.pptx)
    if not path.exists():
        print(f"Error: file not found: {path}", file=sys.stderr)
        sys.exit(1)
    if path.suffix.lower() != ".pptx":
        print(f"Warning: expected a .pptx file, got: {path.suffix}", file=sys.stderr)

    data = extract(path)
    output = json.dumps(data, indent=2, ensure_ascii=False)

    if args.output == "-":
        print(output)
    else:
        out_path = Path(args.output)
        out_path.write_text(output, encoding="utf-8")
        print(f"Saved to {out_path}", file=sys.stderr)


if __name__ == "__main__":
    main()
